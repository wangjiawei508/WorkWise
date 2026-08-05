#!/usr/bin/env python3
"""Native PPTX → SVG converter (bypasses pptx_to_svg approximations).

Reads the PPTX directly with python-pptx and emits one SVG per slide with:
- slide-size coordinates (96 dpi px)
- shapes (rect/ellipse/line/textbox/picture/table/group) with frame geometry
- real text paragraphs: font family/size/weight/color/letter-spacing/alignment
- solid/gradient-first-stop fills, outline colors
- theme color resolution (theme1.xml schemeClr → srgbClr)
- pictures written to `<output>/svg/media/<name>` and referenced as href

WorkWise uses this instead of pptx_to_svg when it yields real editable text.
"""

from __future__ import annotations

import re
import sys
import zipfile
from pathlib import Path
from typing import Any, Optional

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu

EMU_PER_PX = 9525.0
PT_TO_PX = 96.0 / 72.0
XML_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

THEME_KEY_BY_TYPE = {
    "DARK_1": "dk1",
    "DARK_2": "dk2",
    "LIGHT_1": "lt1",
    "LIGHT_2": "lt2",
    "ACCENT_1": "accent1",
    "ACCENT_2": "accent2",
    "ACCENT_3": "accent3",
    "ACCENT_4": "accent4",
    "ACCENT_5": "accent5",
    "ACCENT_6": "accent6",
    "HYPERLINK": "hlink",
    "FOLLOWED_HYPERLINK": "folHlink",
}


def px(value: Any, fallback: float = 0.0) -> float:
    if value is None:
        return fallback
    try:
        return float(value) / EMU_PER_PX
    except (TypeError, ValueError):
        return fallback


def pt_to_px(value: Any, fallback: float = 18.0) -> float:
    if value is None:
        return fallback
    try:
        emu = float(value)
        # python-pptx font.size returns EMU; accept raw pt for small numbers.
        if emu < 1000:
            return emu * PT_TO_PX
        return emu / EMU_PER_PX
    except (TypeError, ValueError):
        return fallback


def esc(value: Any) -> str:
    text = "" if value is None else str(value)
    return (
        text.replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace('"', "&quot;")
    )


def fmt(v: float) -> str:
    return f"{v:.2f}".rstrip("0").rstrip(".")


def resolve_theme_colors(pptx_path: str) -> dict[str, str]:
    colors: dict[str, str] = {}
    try:
        with zipfile.ZipFile(pptx_path) as archive:
            theme_names = [
                name
                for name in archive.namelist()
                if re.match(r"^ppt/theme/theme\d+\.xml$", name)
            ]
            if not theme_names:
                return colors
            root = etree.fromstring(archive.read(theme_names[0]))
            scheme = root.find(f".//{{{XML_NS_A}}}clrScheme")
            if scheme is None:
                return colors
            for child in scheme:
                tag = etree.QName(child).localname
                srgb = child.find(f"{{{XML_NS_A}}}srgbClr")
                if srgb is not None:
                    colors[tag] = "#" + (srgb.get("val") or "").upper()
                    continue
                sysclr = child.find(f"{{{XML_NS_A}}}sysClr")
                if sysclr is not None and sysclr.get("lastClr"):
                    colors[tag] = "#" + sysclr.get("lastClr").upper()
    except Exception:
        pass
    return colors


def resolve_theme_fonts(pptx_path: str) -> dict[str, str]:
    fonts: dict[str, str] = {}
    try:
        with zipfile.ZipFile(pptx_path) as archive:
            theme_names = [
                name
                for name in archive.namelist()
                if re.match(r"^ppt/theme/theme\d+\.xml$", name)
            ]
            if not theme_names:
                return fonts
            root = etree.fromstring(archive.read(theme_names[0]))
            font_scheme = root.find(f".//{{{XML_NS_A}}}fontScheme")
            if font_scheme is None:
                return fonts
            for tag in ("majorFont", "minorFont"):
                node = font_scheme.find(f"{{{XML_NS_A}}}{tag}")
                if node is None:
                    continue
                latin = node.find(f"{{{XML_NS_A}}}latin")
                ea = node.find(f"{{{XML_NS_A}}}ea")
                fonts[tag] = (
                    (latin.get("typeface") if latin is not None else "")
                    or (ea.get("typeface") if ea is not None else "")
                    or "Arial"
                )
    except Exception:
        pass
    return fonts


def color_hex(color: Any, theme_colors: dict[str, str]) -> Optional[str]:
    if color is None:
        return None
    try:
        if color.type is None:
            return None
        type_name = str(color.type)
        if "RGB" in type_name or "NOT_THEME" in type_name:
            rgb: RGBColor = color.rgb
            return "#%02X%02X%02X" % (rgb[0], rgb[1], rgb[2])
        enum_name = type_name.rsplit(".", 1)[-1]
        match = re.match(r"([A-Z][A-Z_0-9]*)", enum_name)
        key = THEME_KEY_BY_TYPE.get(match.group(1) if match else enum_name)
        return theme_colors.get(key) if key else None
    except Exception:
        return None


def shape_fill(shape: Any, theme_colors: dict[str, str]) -> Optional[str]:
    try:
        fill = shape.fill
        if fill.type is None:
            return None
        type_name = str(fill.type)
        if "SOLID" in type_name:
            return color_hex(fill.fore_color, theme_colors) or "#FFFFFF"
        if "GRADIENT" in type_name:
            try:
                stops = list(fill.gradient_stops)
                if stops:
                    return color_hex(stops[0].color, theme_colors) or "#FFFFFF"
            except Exception:
                pass
        return None
    except Exception:
        return None


def shape_stroke(shape: Any, theme_colors: dict[str, str]) -> tuple[Optional[str], Optional[float]]:
    try:
        line = shape.line
        fill = line.fill
        if fill.type is None:
            return None, None
        type_name = str(fill.type)
        color = None
        if "SOLID" in type_name:
            color = color_hex(fill.fore_color, theme_colors)
        width = None
        try:
            if line.width is not None:
                width = max(0.0, px(line.width))
        except Exception:
            width = None
        return color, width
    except Exception:
        return None, None


def run_letter_spacing(run: Any) -> Optional[float]:
    """Return letter spacing in px from the run XML `spc` attribute (1/100 pt)."""
    try:
        rPr = run._r.find(f"{{{XML_NS_A}}}rPr")
        if rPr is None:
            return None
        spc = rPr.get("spc")
        if spc is None:
            return None
        return float(spc) / 100.0 * PT_TO_PX
    except Exception:
        return None


def run_font_name(run: Any, default: str) -> str:
    try:
        rPr = run._r.find(f"{{{XML_NS_A}}}rPr")
        latin = None
        east_asian = None
        if rPr is not None:
            latin_node = rPr.find(f"{{{XML_NS_A}}}latin")
            ea_node = rPr.find(f"{{{XML_NS_A}}}ea")
            latin = latin_node.get("typeface") if latin_node is not None else None
            east_asian = ea_node.get("typeface") if ea_node is not None else None
        if not latin:
            latin = run.font.name or default
        stack = []
        if east_asian and east_asian != latin:
            stack.append(east_asian)
        stack.append(latin or default)
        stack.append("PingFang SC")
        stack.append("Arial")
        return ", ".join(f'"{name}"' for name in stack) + ", sans-serif"
    except Exception:
        return default


def text_anchor(alignment: Any) -> str:
    if alignment == PP_ALIGN.CENTER:
        return "middle"
    if alignment == PP_ALIGN.RIGHT:
        return "end"
    return "start"


def anchor_offset(anchor: Any, content_height: float, box_height: float) -> float:
    if anchor == MSO_ANCHOR.MIDDLE:
        return max(0.0, (box_height - content_height) / 2.0)
    if anchor == MSO_ANCHOR.BOTTOM:
        return max(0.0, box_height - content_height)
    return 0.0


def paragraph_line_height(paragraph: Any, base_font_size_px: float) -> float:
    spacing = paragraph.line_spacing
    if spacing is None:
        return base_font_size_px * 1.2
    try:
        if isinstance(spacing, Emu):
            return max(base_font_size_px, px(spacing))
    except Exception:
        pass
    try:
        return max(base_font_size_px, float(spacing) * base_font_size_px)
    except Exception:
        return base_font_size_px * 1.2


def emit_text(shape: Any, tf: Any, shape_left: float, shape_top: float,
              shape_w: float, shape_h: float, theme_colors: dict[str, str],
              minor_font: str, out: list[str]) -> None:
    try:
        ml = px(tf.margin_left, 7.2)
        mt = px(tf.margin_top, 3.6)
        mb = px(tf.margin_bottom, 3.6)
        anchor = tf.vertical_anchor
    except Exception:
        ml, mt, mb, anchor = 7.2, 3.6, 3.6, None

    paragraphs = list(tf.paragraphs)
    if not paragraphs:
        return

    base_font = minor_font
    # First non-empty run gives the paragraph font baseline.
    first_size = 18.0
    for paragraph in paragraphs:
        for run in paragraph.runs:
            if run.text.strip():
                first_size = pt_to_px(run.font.size, 18.0)
                break
        break

    heights: list[float] = []
    for paragraph in paragraphs:
        if not any(run.text.strip() for run in paragraph.runs):
            heights.append(base_font * 1.2)
            continue
        size = base_font
        for run in paragraph.runs:
            if run.text.strip():
                size = pt_to_px(run.font.size, first_size)
                break
        heights.append(paragraph_line_height(paragraph, size))
    content_height = sum(heights)
    usable_height = max(0.0, shape_h - mt - mb)
    y_offset = mt + anchor_offset(anchor, min(content_height, usable_height), usable_height)

    cursor_y = y_offset
    for paragraph, height in zip(paragraphs, heights):
        text = "".join(run.text for run in paragraph.runs)
        if not text.strip():
            cursor_y += height
            continue
        runs = [run for run in paragraph.runs if run.text.strip()]
        if not runs:
            cursor_y += height
            continue
        first = runs[0]
        font_size = pt_to_px(first.font.size, first_size)
        family = run_font_name(first, base_font)
        bold = bool(first.font.bold)
        italic = bool(first.font.italic)
        fill = color_hex(first.font.color, theme_colors) if first.font.color and first.font.color.type is not None else None
        if not fill:
            fill = shape_fill(shape, theme_colors)
        if not fill:
            fill = "#000000"
        spacing = run_letter_spacing(first)
        align = text_anchor(paragraph.alignment)
        # Baseline: top + y_offset + ascent approximation (0.85em for CJK/Latin mix).
        baseline = shape_top + cursor_y + font_size * 0.85
        if align == "middle":
            x = shape_left + (shape_w - ml - mb) / 2.0 + ml
        elif align == "end":
            x = shape_left + shape_w - mb
        else:
            x = shape_left + ml
        attrs = [
            f'x="{fmt(x)}"',
            f'y="{fmt(baseline)}"',
            f'font-family="{esc(family)}"',
            f'font-size="{fmt(font_size)}"',
            f'fill="{fill}"',
            f'text-anchor="{align}"',
        ]
        if bold:
            attrs.append('font-weight="700"')
        if italic:
            attrs.append('font-style="italic"')
        if spacing:
            attrs.append(f'letter-spacing="{fmt(spacing)}"')
        out.append(f"<text {' '.join(attrs)}>{esc(text)}</text>")
        cursor_y += height


def group_child_context(shape: Any, abs_left: float, abs_top: float,
                        abs_w: float, abs_h: float) -> dict[str, float]:
    ctx: dict[str, float] = {
        "offsetX": abs_left,
        "offsetY": abs_top,
        "scaleX": 1.0,
        "scaleY": 1.0,
        "chOffX": 0.0,
        "chOffY": 0.0,
    }
    try:
        element = shape._element
        grp_sp_pr = element.find("{http://schemas.openxmlformats.org/presentationml/2006/main}grpSpPr")
        xfrm = grp_sp_pr.find(f"{{{XML_NS_A}}}xfrm") if grp_sp_pr is not None else None
        if xfrm is None:
            return ctx
        ch_off = xfrm.find(f"{{{XML_NS_A}}}chOff")
        ch_ext = xfrm.find(f"{{{XML_NS_A}}}chExt")

        def num(node: Any, attr: str) -> float:
            if node is None or node.get(attr) is None:
                return 0.0
            try:
                return float(node.get(attr))
            except (TypeError, ValueError):
                return 0.0

        ch_w = num(ch_ext, "cx")
        ch_h = num(ch_ext, "cy")
        if ch_w > 0 and ch_h > 0 and abs_w > 0 and abs_h > 0:
            ctx["scaleX"] = abs_w / (ch_w / EMU_PER_PX)
            ctx["scaleY"] = abs_h / (ch_h / EMU_PER_PX)
        ctx["chOffX"] = num(ch_off, "x") / EMU_PER_PX
        ctx["chOffY"] = num(ch_off, "y") / EMU_PER_PX
    except Exception:
        pass
    return ctx


TOP_LEVEL_CTX: dict[str, float] = {
    "offsetX": 0.0,
    "offsetY": 0.0,
    "scaleX": 1.0,
    "scaleY": 1.0,
    "chOffX": 0.0,
    "chOffY": 0.0,
}


def shape_to_svg(shape: Any, ctx: dict[str, float],
                 theme_colors: dict[str, str], minor_font: str,
                 media_dir: Path, out: list[str]) -> None:
    try:
        shape_type = shape.shape_type
    except Exception:
        shape_type = None

    # python-pptx reports group-child coordinates in absolute slide space,
    # so no group offset/scale math is applied here.
    left = px(shape.left)
    top = px(shape.top)
    width = max(0.0, px(shape.width))
    height = max(0.0, px(shape.height))
    rotation = 0.0
    try:
        rotation = float(shape.rotation or 0.0)
    except Exception:
        rotation = 0.0

    fill = shape_fill(shape, theme_colors)
    stroke, stroke_width = shape_stroke(shape, theme_colors)

    transform = ""
    if rotation:
        cx = left + width / 2.0
        cy = top + height / 2.0
        transform = f' transform="rotate({fmt(rotation)} {fmt(cx)} {fmt(cy)})"'

    common = (
        f'x="{fmt(left)}" y="{fmt(top)}" width="{fmt(max(width, 0.1))}" '
        f'height="{fmt(max(height, 0.1))}"'
    )
    style = []
    if fill:
        style.append(f'fill="{fill}"')
    else:
        style.append('fill="none"')
    if stroke:
        style.append(f'stroke="{stroke}"')
        style.append(f'stroke-width="{fmt(stroke_width if stroke_width is not None else 1.0)}"')
    else:
        style.append('stroke="none"')

    if shape_type == MSO_SHAPE_TYPE.GROUP:
        try:
            for child in shape.shapes:
                shape_to_svg(child, TOP_LEVEL_CTX, theme_colors, minor_font, media_dir, out)
        except Exception:
            pass
        return

    if shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            image = shape.image
            ext = image.ext or "png"
            filename = f"pic_{shape.shape_id}.{ext}"
            (media_dir / filename).write_bytes(image.blob)
            out.append(f'<image href="media/{filename}" {common}{transform}/>')
        except Exception:
            out.append(f'<rect {common} fill="#D9D9D9" stroke="#B0B0B0" stroke-width="1"{transform}/>')
        return

    if shape_type == MSO_SHAPE_TYPE.LINE:
        out.append(
            f'<line x1="{fmt(left)}" y1="{fmt(top)}" x2="{fmt(left + width)}" '
            f'y2="{fmt(top + height)}" stroke="{stroke or "#000000"}" '
            f'stroke-width="{fmt(stroke_width if stroke_width is not None else 1.0)}"{transform}/>'
        )
        return

    if shape_type == MSO_SHAPE_TYPE.TABLE:
        try:
            table = shape.table
            row_heights = [px(row.height) for row in table.rows]
            col_widths = [px(col.width) for col in table.columns]
            for row_index, row in enumerate(table.rows):
                cell_y = top + sum(row_heights[:row_index])
                cell_x = left
                for col_index, cell in enumerate(row.cells):
                    cell_w = col_widths[col_index]
                    cell_h = row_heights[row_index]
                    out.append(
                        f'<rect x="{fmt(cell_x)}" y="{fmt(cell_y)}" '
                        f'width="{fmt(cell_w)}" height="{fmt(cell_h)}" fill="{fill or "#FFFFFF"}" '
                        f'stroke="#D1D5DB" stroke-width="1"/>'
                    )
                    emit_text(shape, cell.text_frame, cell_x, cell_y, cell_w, cell_h,
                              theme_colors, minor_font, out)
                    cell_x += cell_w
        except Exception:
            out.append(f'<rect {common} fill="{fill or "#FFFFFF"}" stroke="#D1D5DB" stroke-width="1"{transform}/>')
        return

    try:
        is_ellipse = shape.auto_shape_type is not None and "OVAL" in str(shape.auto_shape_type).upper()
    except Exception:
        is_ellipse = False
    if is_ellipse:
        out.append(
            f'<ellipse cx="{fmt(left + width / 2.0)}" cy="{fmt(top + height / 2.0)}" '
            f'rx="{fmt(width / 2.0)}" ry="{fmt(height / 2.0)}" {" ".join(style)}{transform}/>'
        )
    else:
        out.append(f'<rect {common} {" ".join(style)}{transform}/>')

    # Text frame
    try:
        if shape.has_text_frame:
            emit_text(shape, shape.text_frame, left, top, width, height, theme_colors, minor_font, out)
    except Exception:
        pass


def convert_pptx_to_svg(pptx_path: str, output_dir: str) -> dict[str, Any]:
    output = Path(output_dir)
    svg_dir = output / "svg"
    media_dir = svg_dir / "media"
    svg_dir.mkdir(parents=True, exist_ok=True)
    media_dir.mkdir(parents=True, exist_ok=True)

    theme_colors = resolve_theme_colors(pptx_path)
    theme_fonts = resolve_theme_fonts(pptx_path)
    minor_font = theme_fonts.get("minorFont") or theme_fonts.get("majorFont") or "Arial"

    prs = Presentation(pptx_path)
    slide_w = px(prs.slide_width, 1280.0)
    slide_h = px(prs.slide_height, 720.0)

    converted = 0
    for index, slide in enumerate(prs.slides, start=1):
        out: list[str] = [
            f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 {fmt(slide_w)} {fmt(slide_h)}" '
            f'width="{fmt(slide_w)}" height="{fmt(slide_h)}">'
        ]
        background = "#FFFFFF"
        try:
            bg_fill = shape_fill(slide.background, theme_colors)
            if bg_fill:
                background = bg_fill
        except Exception:
            pass
        out.append(f'<rect x="0" y="0" width="{fmt(slide_w)}" height="{fmt(slide_h)}" fill="{background}"/>')
        for shape in slide.shapes:
            shape_to_svg(shape, TOP_LEVEL_CTX, theme_colors, minor_font, media_dir, out)
        out.append("</svg>")
        (svg_dir / f"slide_{index:02d}.svg").write_text("\n".join(out), encoding="utf-8")
        converted += 1

    return {"slides": converted, "mediaFiles": len(list(media_dir.iterdir()))}


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("usage: pptx_native_to_svg.py <input.pptx> <output_dir>", file=sys.stderr)
        sys.exit(2)
    result = convert_pptx_to_svg(sys.argv[1], sys.argv[2])
    print(f"Converted {result['slides']} slides, {result['mediaFiles']} media files.")
