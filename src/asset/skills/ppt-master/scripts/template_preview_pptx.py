#!/usr/bin/env python3
"""
PPT Master - Template Preview PPTX Exporter

Export public SVG prototypes as a structured review deck while retaining
definition-only Layout prototypes in the native package.

Usage:
    python3 scripts/template_preview_pptx.py <template_workspace> [-o output.pptx]

Examples:
    python3 scripts/template_preview_pptx.py projects/my_template
    python3 scripts/template_preview_pptx.py templates/decks/my_template -o review.pptx
    python3 scripts/template_preview_pptx.py templates/decks/legacy --visual-only

Dependencies:
    python-pptx
"""

from __future__ import annotations

import argparse
import contextlib
import math
import re
import shutil
import statistics
import sys
import tempfile
from collections.abc import Iterator
from pathlib import Path
from xml.etree import ElementTree as ET

from console_encoding import configure_utf8_stdio


configure_utf8_stdio()

from pptx import Presentation  # noqa: E402

from svg_to_pptx.drawingml.theme_fonts import (  # noqa: E402
    MasterTextStyleSpec,
)
from svg_to_pptx.drawingml.utils import font_px_to_hpt  # noqa: E402
from svg_to_pptx.pptx_package.builder import (  # noqa: E402
    create_pptx_with_native_svg,
)


_FRONTMATTER_ID_RE = re.compile(
    r"^(?:template_id|deck_id|layout_id)\s*:\s*(.+?)\s*$",
    re.MULTILINE,
)
_REPLICATION_MODE_RE = re.compile(
    r"^replication_mode\s*:\s*(standard|fidelity|mirror)\s*$",
    re.MULTILINE,
)
_CANVAS_VIEWBOX_RE = re.compile(
    r"^canvas_viewbox\s*:\s*[\"']?([^\"'\r\n]+?)[\"']?\s*$",
    re.MULTILINE,
)
_FONT_SIZE_RE = re.compile(r"^([0-9]+(?:\.[0-9]+)?)(?:px)?$")
_FILENAME_UNSAFE_RE = re.compile(r"[\\/:*?\"<>|\x00-\x1f]+")
_PLACEHOLDER_MARKER_RE = re.compile(r"\{\{([A-Z][A-Z0-9_]*)\}\}")
_TITLE_PLACEHOLDERS = frozenset({"title", "subtitle"})
_BODY_PLACEHOLDERS = frozenset({
    "body",
    "date",
    "footer",
    "slide-number",
})
_DEFAULT_TITLE_PX = 40.0
_DEFAULT_BODY_PX = 24.0


def _review_marker_text(match: re.Match[str]) -> str:
    """Return concise preview-only text for one canonical marker."""
    token = match.group(1)
    if token in {"PAGE_NUM", "SLIDE_NUM"}:
        return "1"
    if token.endswith("_NUM"):
        return "01"
    if token == "DATE":
        return "YYYY-MM-DD"
    return token.replace("_", " ").title()


def _write_review_svg(source: Path, target: Path) -> bool:
    """Copy one SVG, shortening only visible placeholder-carrier prompts."""
    tree = ET.parse(source)
    changed = False
    for slot in tree.getroot().iter():
        if not (slot.get("data-pptx-placeholder") or "").strip():
            continue
        for carrier in slot.iter():
            if (
                carrier.get("data-pptx-carrier") or ""
            ).strip().lower() != "true":
                continue
            for element in carrier.iter():
                if element.text:
                    updated = _PLACEHOLDER_MARKER_RE.sub(
                        _review_marker_text,
                        element.text,
                    )
                    if updated != element.text:
                        element.text = updated
                        changed = True
    if changed:
        tree.write(target, encoding="utf-8", xml_declaration=True)
    else:
        shutil.copy2(source, target)
    return changed


@contextlib.contextmanager
def _review_svg_sources(
    workspace: Path,
    svg_files: list[Path],
    *,
    shorten_placeholder_markers: bool,
) -> Iterator[list[Path]]:
    """Yield ephemeral review SVGs without modifying canonical template files."""
    if not shorten_placeholder_markers:
        yield svg_files
        return

    with tempfile.TemporaryDirectory(
        prefix=".template-preview-",
        dir=workspace,
    ) as temporary:
        review_dir = Path(temporary)
        review_files: list[Path] = []
        shortened = 0
        for source in svg_files:
            target = review_dir / source.name
            shortened += int(_write_review_svg(source, target))
            review_files.append(target)
        print(
            "  Review prompt text: preview-only samples in "
            f"{shortened} SVG(s); canonical {{{{...}}}} markers unchanged"
        )
        yield review_files


def _partition_svg_prototypes(
    svg_files: list[Path],
    *,
    visual_only: bool,
) -> tuple[list[Path], list[Path]]:
    """Separate public pages from canonical definition-only Layout SVGs."""
    if visual_only:
        return svg_files, []
    public_files: list[Path] = []
    definition_files: list[Path] = []
    for path in svg_files:
        target = definition_files if path.stem.startswith("layout_") else public_files
        target.append(path)
    return public_files, definition_files


def _resolve_workspace(path: Path) -> tuple[Path, Path]:
    """Resolve one workspace root and its canonical template-source directory."""
    candidate = path.expanduser().resolve()
    nested_spec = candidate / "templates" / "design_spec.md"
    if nested_spec.is_file():
        return candidate, candidate / "templates"

    direct_spec = candidate / "design_spec.md"
    if direct_spec.is_file():
        if candidate.name == "templates" and (candidate.parent / "exports").is_dir():
            return candidate.parent, candidate
        return candidate, candidate

    raise ValueError(
        "template workspace must contain templates/design_spec.md "
        "(current structure) or design_spec.md (legacy flat package)"
    )


def _template_id(spec_path: Path, workspace: Path) -> str:
    """Read a portable template id, falling back to the workspace directory name."""
    text = spec_path.read_text(encoding="utf-8")
    match = _FRONTMATTER_ID_RE.search(text)
    raw = match.group(1).strip().strip("'\"") if match else workspace.name
    safe = _FILENAME_UNSAFE_RE.sub("_", raw).strip(" ._")
    return safe or "template"


def _replication_mode(spec_path: Path) -> str:
    """Read the template replication mode, defaulting legacy packages to standard."""
    text = spec_path.read_text(encoding="utf-8")
    match = _REPLICATION_MODE_RE.search(text)
    return match.group(1) if match else "standard"


def _canvas_viewbox(spec_path: Path) -> str | None:
    """Read the template's locked root canvas when declared."""
    text = spec_path.read_text(encoding="utf-8")
    if not text.startswith("---\n"):
        return None
    end = text.find("\n---\n", 4)
    if end == -1:
        return None
    match = _CANVAS_VIEWBOX_RE.search(text[4:end])
    return match.group(1).strip() if match else None


def _style_property(style: str, name: str) -> str | None:
    """Return one inline CSS declaration value."""
    for declaration in style.split(";"):
        key, separator, value = declaration.partition(":")
        if separator and key.strip().lower() == name:
            return value.strip()
    return None


def _font_size_px(element: ET.Element) -> float | None:
    """Read one finite positive SVG font size in px."""
    raw = element.get("font-size")
    if raw is None:
        raw = _style_property(element.get("style", ""), "font-size")
    if raw is None:
        return None
    match = _FONT_SIZE_RE.fullmatch(raw.strip())
    if match is None:
        return None
    value = float(match.group(1))
    return value if math.isfinite(value) and value > 0 else None


def _carrier_sizes(svg_files: list[Path]) -> tuple[list[float], list[float]]:
    """Collect authored title/body sizes from semantic placeholder carriers."""
    title_sizes: list[float] = []
    body_sizes: list[float] = []
    for svg_path in svg_files:
        root = ET.parse(svg_path).getroot()
        for slot in root.iter():
            placeholder = slot.get("data-pptx-placeholder")
            if placeholder not in _TITLE_PLACEHOLDERS | _BODY_PLACEHOLDERS:
                continue
            for carrier in slot.iter():
                if carrier.get("data-pptx-carrier") != "true":
                    continue
                size = _font_size_px(carrier)
                if size is None:
                    continue
                target = title_sizes if placeholder in _TITLE_PLACEHOLDERS else body_sizes
                target.append(size)
    return title_sizes, body_sizes


def _master_text_style(svg_files: list[Path]) -> tuple[MasterTextStyleSpec, float, float]:
    """Build review-only Master text defaults without requiring a project lock."""
    title_sizes, body_sizes = _carrier_sizes(svg_files)
    title_px = float(statistics.median(title_sizes)) if title_sizes else _DEFAULT_TITLE_PX
    body_px = float(statistics.median(body_sizes)) if body_sizes else _DEFAULT_BODY_PX
    return (
        MasterTextStyleSpec(
            title_hpt=font_px_to_hpt(title_px),
            body_hpt=font_px_to_hpt(body_px),
        ),
        title_px,
        body_px,
    )


def _verify_output(
    output_path: Path,
    *,
    require_full_placeholder_frames: bool,
) -> tuple[int, int, int, int]:
    """Reopen the review deck and verify counts plus authored placeholder frames."""
    presentation = Presentation(str(output_path))
    master_count = len(presentation.slide_masters)
    layout_count = sum(len(master.slide_layouts) for master in presentation.slide_masters)
    placeholder_count = 0
    if require_full_placeholder_frames:
        for slide_number, slide in enumerate(presentation.slides, 1):
            layout_placeholders = {
                shape.placeholder_format.idx: shape
                for shape in slide.slide_layout.placeholders
            }
            slide_placeholders = {
                shape.placeholder_format.idx: shape
                for shape in slide.placeholders
            }
            if set(slide_placeholders) != set(layout_placeholders):
                raise ValueError(
                    f"review slide {slide_number} placeholder indexes do not match "
                    f"its Layout: {sorted(slide_placeholders)} != "
                    f"{sorted(layout_placeholders)}"
                )
            for placeholder_idx, slide_shape in slide_placeholders.items():
                layout_shape = layout_placeholders[placeholder_idx]
                if (
                    slide_shape.placeholder_format.type
                    != layout_shape.placeholder_format.type
                ):
                    raise ValueError(
                        f"review slide {slide_number} placeholder {placeholder_idx} "
                        "type does not match its Layout"
                    )
                slide_frame = (
                    slide_shape.left,
                    slide_shape.top,
                    slide_shape.width,
                    slide_shape.height,
                )
                layout_frame = (
                    layout_shape.left,
                    layout_shape.top,
                    layout_shape.width,
                    layout_shape.height,
                )
                if slide_frame != layout_frame:
                    raise ValueError(
                        f"review slide {slide_number} placeholder {placeholder_idx} "
                        f"uses a tight/local frame {slide_frame}; expected full "
                        f"Layout frame {layout_frame}"
                    )
                placeholder_count += 1
    return len(presentation.slides), master_count, layout_count, placeholder_count


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description=(
            "Export a complete template workspace as a structured PPTX review deck."
        ),
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    parser.add_argument(
        "template_workspace",
        help=(
            "Workspace containing templates/design_spec.md; legacy flat template "
            "directories are also accepted."
        ),
    )
    parser.add_argument(
        "-o",
        "--output",
        help=(
            "Output PPTX path. Default: "
            "<template_workspace>/exports/<template_id>_template_preview.pptx"
        ),
    )
    parser.add_argument(
        "--force",
        action="store_true",
        help="Replace an existing review PPTX after an intentional re-export.",
    )
    parser.add_argument(
        "--visual-only",
        action="store_true",
        help=(
            "Export a legacy SVG roster as slide-local DrawingML for visual review. "
            "This does not validate or claim a reusable Master/Layout contract."
        ),
    )
    return parser


def main(argv: list[str] | None = None) -> int:
    parser = build_parser()
    args = parser.parse_args(argv)

    try:
        workspace, template_dir = _resolve_workspace(Path(args.template_workspace))
        all_svg_files = sorted(template_dir.glob("*.svg"))
        if not all_svg_files:
            raise ValueError(f"template directory has no SVG prototypes: {template_dir}")
        svg_files, layout_definition_files = _partition_svg_prototypes(
            all_svg_files,
            visual_only=args.visual_only,
        )
        if not svg_files:
            raise ValueError(
                "template directory contains Layout definitions but no public "
                f"SVG prototypes: {template_dir}"
            )

        spec_path = template_dir / "design_spec.md"
        template_id = _template_id(spec_path, workspace)
        replication_mode = _replication_mode(spec_path)
        locked_canvas = _canvas_viewbox(spec_path)
        if locked_canvas is None and not args.visual_only:
            raise ValueError(
                "design_spec.md frontmatter must declare canvas_viewbox"
            )
        use_full_placeholder_frames = (
            not args.visual_only and replication_mode != "mirror"
        )
        output_path = (
            Path(args.output).expanduser().resolve()
            if args.output
            else workspace / "exports" / f"{template_id}_template_preview.pptx"
        )
        if output_path.suffix.lower() != ".pptx":
            raise ValueError(f"output must use a .pptx extension: {output_path}")
        if output_path.exists() and not args.force:
            raise ValueError(
                f"output already exists: {output_path}; use --force to replace it"
            )
        output_path.parent.mkdir(parents=True, exist_ok=True)
        text_style: MasterTextStyleSpec | None = None
        if not args.visual_only:
            text_style, title_px, body_px = _master_text_style(all_svg_files)

        print("PPT Master - Template Preview PPTX Exporter")
        print(f"  Workspace: {workspace}")
        print(f"  Template source: {template_dir}")
        print(f"  Public SVG prototypes: {len(svg_files)}")
        if layout_definition_files:
            print(
                "  Definition-only Layout prototypes: "
                f"{len(layout_definition_files)}"
            )
        if args.visual_only:
            print("  Review mode: visual-only legacy compatibility")
        elif replication_mode == "mirror":
            print("  Review placeholder frames: preserved source Slide geometry")
        else:
            print(f"  Review Master defaults: title {title_px:g}px, body {body_px:g}px")
            print("  Review placeholder frames: full Layout bounds")
        print(f"  Output: {output_path}")

        with _review_svg_sources(
            workspace,
            all_svg_files,
            shorten_placeholder_markers=use_full_placeholder_frames,
        ) as review_all_svg_files:
            review_svg_files, review_layout_definition_files = (
                _partition_svg_prototypes(
                    review_all_svg_files,
                    visual_only=args.visual_only,
                )
            )
            success = create_pptx_with_native_svg(
                svg_files=review_svg_files,
                output_path=output_path,
                canvas_format=None,
                expected_viewbox=locked_canvas,
                verbose=True,
                transition=None,
                enable_notes=False,
                animation=None,
                image_optimize=False,
                native_objects=True,
                pptx_structure="flat" if args.visual_only else "structured",
                use_layout_placeholder_frames=use_full_placeholder_frames,
                master_text_style_spec=text_style,
                structure_name=template_id,
                layout_definition_files=review_layout_definition_files,
            )
        if not success or not output_path.is_file():
            print("Error: template preview export did not produce a PPTX", file=sys.stderr)
            return 1

        slide_count, master_count, layout_count, placeholder_count = _verify_output(
            output_path,
            require_full_placeholder_frames=use_full_placeholder_frames,
        )
        if slide_count != len(svg_files):
            print(
                "Error: review PPTX slide count does not match the template SVG roster "
                f"({slide_count} != {len(svg_files)})",
                file=sys.stderr,
            )
            return 1

        label = "Visual-only template preview" if args.visual_only else "Template preview"
        placeholder_status = (
            f", {placeholder_count} full-frame placeholder(s)"
            if use_full_placeholder_frames
            else ""
        )
        print(
            f"[OK] {label} verified: "
            f"{slide_count} slides, {master_count} master(s), "
            f"{layout_count} layout(s){placeholder_status}"
        )
        print(output_path)
        return 0
    except (OSError, ET.ParseError, RuntimeError, ValueError) as exc:
        print(f"Error: {exc}", file=sys.stderr)
        return 1


if __name__ == "__main__":
    raise SystemExit(main())
